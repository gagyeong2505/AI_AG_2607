"""
스크랩한 기사 콘텐츠 JSON을 받아 PowerPoint(.pptx) 발표자료를 생성한다.
표지 -> 개요 -> 기사별 슬라이드(스크린샷 좌측 / 요약·핵심포인트 우측) -> 결론 순서로 구성하며
맑은 고딕과 설정 가능한 메인 컬러를 일관되게 적용한다.

입력 JSON 형식은 ../references/content_schema.md 참고 (generate_docx.py와 동일한 파일을 공유한다).

사용법:
    python generate_pptx.py --input content.json --output deck.pptx [--color 2E74B5]
"""

import argparse
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FONT_NAME = "맑은 고딕"
DEFAULT_MAIN_COLOR = "2E74B5"

TITLE_LAYOUT = 0
BULLETS_LAYOUT = 1
TITLE_ONLY_LAYOUT = 5

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def clear_existing_slides(prs: Presentation) -> None:
    """템플릿 pptx에 들어있던 샘플 슬라이드를 모두 제거하고 마스터/레이아웃/테마만 남긴다."""
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        xml_slides.remove(slide_id)


def hex_to_rgbcolor(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        raise ValueError(f"올바른 hex 색상 코드가 아닙니다: {hex_color}")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def style_run(run, size_pt: int, bold: bool = False, color: RGBColor = None, italic: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def set_title_text(slide, text: str, size_pt: int, color: RGBColor) -> None:
    title_shape = slide.shapes.title
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            style_run(run, size_pt, bold=True, color=color)


def add_title_slide(prs: Presentation, title: str, subtitle: str, color: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    set_title_text(slide, title, 40, color)
    if subtitle and len(slide.placeholders) > 1:
        sub_shape = slide.placeholders[1]
        sub_shape.text = subtitle
        for paragraph in sub_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                style_run(run, 20, italic=True)
    return slide


def add_bullets_slide(prs: Presentation, heading: str, bullets: list, color: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[BULLETS_LAYOUT])
    set_title_text(slide, heading, 32, color)
    body = slide.placeholders[1]
    text_frame = body.text_frame
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        for run in paragraph.runs:
            style_run(run, 20)
    return slide


def add_article_slide(prs: Presentation, article: dict, color: RGBColor, slide_width: int, slide_height: int, scale: float = 1.0):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_LAYOUT])
    set_title_text(slide, article["headline"], 28, color)

    content_top = Emu(int(Inches(1.3) * scale))
    left_margin = Emu(int(Inches(0.4) * scale))
    image_width = Emu(int(Inches(5.9) * scale))

    screenshot = article.get("screenshot")
    if screenshot and os.path.exists(screenshot):
        slide.shapes.add_picture(screenshot, left_margin, content_top, width=image_width)
        text_left = Emu(int(Inches(6.6) * scale))
        text_width = slide_width - text_left - Emu(int(Inches(0.4) * scale))
    else:
        if screenshot:
            print(f"경고: 스크린샷을 찾을 수 없어 건너뜁니다: {screenshot}")
        text_left = left_margin
        text_width = slide_width - Emu(int(Inches(0.8) * scale))

    text_box = slide.shapes.add_textbox(text_left, content_top, text_width, slide_height - content_top - Emu(int(Inches(0.4) * scale)))
    tf = text_box.text_frame
    tf.word_wrap = True

    meta_bits = [article["source"]]
    if article.get("published"):
        meta_bits.append(article["published"])
    meta_paragraph = tf.paragraphs[0]
    meta_paragraph.text = " · ".join(meta_bits)
    for run in meta_paragraph.runs:
        style_run(run, 14, italic=True, color=RGBColor(0x60, 0x60, 0x60))

    url_paragraph = tf.add_paragraph()
    url_paragraph.text = article["url"]
    for run in url_paragraph.runs:
        style_run(run, 12, italic=True, color=RGBColor(0x80, 0x80, 0x80))

    summary_paragraph = tf.add_paragraph()
    summary_paragraph.text = article["summary"]
    summary_paragraph.space_before = Pt(12)
    for run in summary_paragraph.runs:
        style_run(run, 18)

    for point in article.get("key_points", []):
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.space_before = Pt(6)
        for run in p.runs:
            style_run(run, 16)

    return slide


def build_deck(content: dict, output_path: str, main_color: str, template_path: str = None) -> None:
    color = hex_to_rgbcolor(main_color)

    if template_path:
        prs = Presentation(template_path)
        clear_existing_slides(prs)
        # 템플릿의 슬라이드 크기(테마 캔버스)는 그대로 유지하고, 기본 13.333x7.5in 기준
        # 좌표를 비율에 맞게 스케일링해 기사 슬라이드 배치가 캔버스를 벗어나지 않게 한다.
        scale = prs.slide_width / SLIDE_WIDTH
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        scale = 1.0

    add_title_slide(prs, content["title"], content.get("subtitle", ""), color)

    intro = content.get("intro", [])
    if intro:
        add_bullets_slide(prs, "개요", intro, color)

    for article in content.get("articles", []):
        add_article_slide(prs, article, color, prs.slide_width, prs.slide_height, scale)

    conclusion = content.get("conclusion", [])
    if conclusion:
        add_bullets_slide(prs, "결론", conclusion, color)

    prs.save(output_path)
    print(f"saved: {output_path}")


def main() -> None:
    parser = argparse.ArgumentParser(description="스크랩 기사 콘텐츠 JSON으로 pptx 발표자료 생성")
    parser.add_argument("--input", required=True, help="콘텐츠 JSON 파일 경로")
    parser.add_argument("--output", required=True, help="저장할 .pptx 파일 경로")
    parser.add_argument("--color", default=DEFAULT_MAIN_COLOR, help="메인 컬러 hex 코드 (# 없이, 예: 2E74B5)")
    parser.add_argument("--template", default=None, help="테마/마스터로 사용할 .pptx 템플릿 경로 (기존 샘플 슬라이드는 제거하고 마스터만 재사용)")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        content = json.load(f)

    output_dir = os.path.dirname(args.output)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    build_deck(content, args.output, args.color, args.template)


if __name__ == "__main__":
    main()
